# src/chuk_mcp_linkedin/utils/document_converter.py
"""
Document to image converter for LinkedIn post previews.

Converts PDF, PowerPoint, and Word documents to images for preview rendering.
"""

from pathlib import Path
from typing import List, Optional
import hashlib


class DocumentConverter:
    """Convert documents (PDF, PPTX, DOCX) to images for preview"""

    # Cache directory for converted images
    CACHE_DIR = Path.home() / ".linkedin_drafts" / "document_cache"

    @staticmethod
    def _get_cache_key(filepath: str) -> str:
        """Generate cache key based on file path and modification time"""
        path = Path(filepath)
        if not path.exists():
            return ""

        # Use file path + modification time for cache key
        mtime = path.stat().st_mtime
        cache_input = f"{filepath}_{mtime}"
        return hashlib.md5(cache_input.encode()).hexdigest()

    @staticmethod
    def _get_cache_dir(cache_key: str) -> Path:
        """Get cache directory for a specific document"""
        cache_dir = DocumentConverter.CACHE_DIR / cache_key
        cache_dir.mkdir(parents=True, exist_ok=True)
        return cache_dir

    @staticmethod
    def convert_to_images(
        filepath: str,
        max_pages: Optional[int] = None,
        dpi: int = 150,
    ) -> List[str]:
        """
        Convert document to images.

        Args:
            filepath: Path to document file
            max_pages: Maximum number of pages to convert (None for all)
            dpi: DPI for image conversion (higher = better quality, larger files)

        Returns:
            List of paths to generated images
        """
        path = Path(filepath)
        if not path.exists():
            raise FileNotFoundError(f"Document not found: {filepath}")

        # Check cache first
        cache_key = DocumentConverter._get_cache_key(filepath)
        if cache_key:
            cached_images = DocumentConverter._get_cached_images(cache_key, max_pages)
            if cached_images:
                return cached_images

        # Convert based on file type
        file_ext = path.suffix.lower()

        if file_ext == ".pdf":
            images = DocumentConverter._convert_pdf(filepath, max_pages, dpi, cache_key)
        elif file_ext in [".pptx", ".ppt"]:
            images = DocumentConverter._convert_pptx(filepath, max_pages, dpi, cache_key)
        elif file_ext in [".docx", ".doc"]:
            images = DocumentConverter._convert_docx(filepath, max_pages, dpi, cache_key)
        else:
            raise ValueError(f"Unsupported file type: {file_ext}")

        return images

    @staticmethod
    def _get_cached_images(cache_key: str, max_pages: Optional[int] = None) -> List[str]:
        """Get cached images if they exist"""
        cache_dir = DocumentConverter.CACHE_DIR / cache_key
        if not cache_dir.exists():
            return []

        # Find all PNG files in cache dir, sorted by name
        images = sorted(cache_dir.glob("page_*.png"))
        image_paths = [str(img) for img in images]

        if max_pages:
            image_paths = image_paths[:max_pages]

        # Return only if we have images
        return image_paths if image_paths else []

    @staticmethod
    def _convert_pdf(
        filepath: str,
        max_pages: Optional[int],
        dpi: int,
        cache_key: str,
    ) -> List[str]:
        """Convert PDF to images using pdf2image"""
        try:
            from pdf2image import convert_from_path
        except ImportError:
            raise ImportError(
                "pdf2image is required for PDF conversion. "
                "Install with: pip install pdf2image\n"
                "Note: pdf2image also requires poppler to be installed:\n"
                "  macOS: brew install poppler\n"
                "  Ubuntu: sudo apt-get install poppler-utils\n"
                "  Windows: Download from https://github.com/oschwartz10612/poppler-windows/releases"
            )

        cache_dir = DocumentConverter._get_cache_dir(cache_key)

        # Convert PDF pages to images
        try:
            images = convert_from_path(
                filepath,
                dpi=dpi,
                fmt="png",
                thread_count=4,
                first_page=1,
                last_page=max_pages if max_pages else None,
            )
        except Exception as e:
            raise RuntimeError(f"Failed to convert PDF: {str(e)}")

        # Save images to cache
        image_paths = []
        for i, img in enumerate(images, start=1):
            output_path = cache_dir / f"page_{i:03d}.png"
            img.save(output_path, "PNG", optimize=True)
            image_paths.append(str(output_path))

        return image_paths

    @staticmethod
    def _convert_pptx(
        filepath: str,
        max_pages: Optional[int],
        dpi: int,
        cache_key: str,
    ) -> List[str]:
        """Convert PowerPoint to images"""
        try:
            from pptx import Presentation
            from PIL import Image
        except ImportError:
            raise ImportError(
                "python-pptx and Pillow are required for PowerPoint conversion. "
                "Install with: pip install python-pptx Pillow"
            )

        cache_dir = DocumentConverter._get_cache_dir(cache_key)

        # Load presentation
        prs = Presentation(filepath)
        slides = prs.slides

        if max_pages:
            slides = list(slides)[:max_pages]

        image_paths = []

        # PowerPoint to image conversion is complex - we'll use a workaround
        # For now, we'll create placeholder images with slide content
        # In production, you might want to use LibreOffice or similar for conversion

        for i, slide in enumerate(slides, start=1):
            output_path = cache_dir / f"page_{i:03d}.png"

            # Create a simple image representation
            # This is a simplified approach - for production use LibreOffice conversion
            width = int(prs.slide_width * dpi / 914400)  # EMUs to pixels
            height = int(prs.slide_height * dpi / 914400)

            # Create blank canvas
            img = Image.new("RGB", (width, height), "white")

            # TODO: Render slide content properly
            # For now, save blank slide representation
            # In production, use:
            # - LibreOffice headless conversion: soffice --headless --convert-to pdf --outdir /tmp file.pptx
            # - Then convert PDF to images using pdf2image

            img.save(output_path, "PNG", optimize=True)
            image_paths.append(str(output_path))

        return image_paths

    @staticmethod
    def _convert_docx(
        filepath: str,
        max_pages: Optional[int],
        dpi: int,
        cache_key: str,
    ) -> List[str]:
        """Convert Word document to images"""
        try:
            from docx import Document
            from PIL import Image
        except ImportError:
            raise ImportError(
                "python-docx and Pillow are required for Word conversion. "
                "Install with: pip install python-docx Pillow"
            )

        cache_dir = DocumentConverter._get_cache_dir(cache_key)

        # Load document
        doc = Document(filepath)

        # Word documents don't have "pages" in the same way
        # We'll estimate based on content or convert to PDF first
        # For now, create placeholder images

        # Estimate page count based on content length (rough approximation)
        total_chars = sum(len(p.text) for p in doc.paragraphs)
        estimated_pages = max(1, total_chars // 3000)  # ~3000 chars per page

        if max_pages:
            estimated_pages = min(estimated_pages, max_pages)

        image_paths = []

        # Create placeholder images
        # In production, convert DOCX to PDF first, then use pdf2image
        for i in range(1, estimated_pages + 1):
            output_path = cache_dir / f"page_{i:03d}.png"

            # Create blank page
            width, height = 816, 1056  # A4 at 96 DPI
            img = Image.new("RGB", (width, height), "white")

            # TODO: Render document content properly
            # For production use:
            # - LibreOffice headless conversion: soffice --headless --convert-to pdf --outdir /tmp file.docx
            # - Then convert PDF to images using pdf2image

            img.save(output_path, "PNG", optimize=True)
            image_paths.append(str(output_path))

        return image_paths

    @staticmethod
    def get_page_count(filepath: str) -> int:
        """Get number of pages in document"""
        path = Path(filepath)
        if not path.exists():
            return 0

        file_ext = path.suffix.lower()

        try:
            if file_ext == ".pdf":
                try:
                    import PyPDF2

                    with open(filepath, "rb") as f:
                        pdf = PyPDF2.PdfReader(f)
                        return len(pdf.pages)
                except ImportError:
                    # Fallback: try pdf2image
                    try:
                        from pdf2image import pdfinfo_from_path

                        info = pdfinfo_from_path(filepath)
                        return info.get("Pages", 0)
                    except ImportError:
                        return 0

            elif file_ext in [".pptx", ".ppt"]:
                from pptx import Presentation

                prs = Presentation(filepath)
                return len(prs.slides)

            elif file_ext in [".docx", ".doc"]:
                # Word documents don't have clear page boundaries
                # Return estimated count
                from docx import Document

                doc = Document(filepath)
                total_chars = sum(len(p.text) for p in doc.paragraphs)
                return max(1, total_chars // 3000)

        except Exception:
            return 0

        return 0

    @staticmethod
    def clear_cache(cache_key: Optional[str] = None):
        """Clear document cache"""
        if cache_key:
            # Clear specific document cache
            cache_dir = DocumentConverter.CACHE_DIR / cache_key
            if cache_dir.exists():
                import shutil

                shutil.rmtree(cache_dir)
        else:
            # Clear all cache
            if DocumentConverter.CACHE_DIR.exists():
                import shutil

                shutil.rmtree(DocumentConverter.CACHE_DIR)
