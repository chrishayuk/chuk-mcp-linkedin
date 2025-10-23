"""
Create dummy test files (PDF, PPTX, DOC) for LinkedIn document upload testing.

These are minimal valid files that LinkedIn should accept.
"""

from pathlib import Path

# Try to create a simple PDF
try:
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas

    def create_pdf():
        pdf_path = Path(__file__).parent / "test_document.pdf"
        c = canvas.Canvas(str(pdf_path), pagesize=letter)

        # Add content
        c.setFont("Helvetica-Bold", 24)
        c.drawString(100, 750, "Test LinkedIn Document")

        c.setFont("Helvetica", 14)
        c.drawString(100, 700, "This is a test PDF for LinkedIn upload.")
        c.drawString(100, 670, "Created by chuk-mcp-linkedin")

        # Add some pages
        for i in range(2, 4):
            c.showPage()
            c.setFont("Helvetica", 16)
            c.drawString(100, 750, f"Page {i}")
            c.drawString(100, 700, "More content here...")

        c.save()
        print(f"✓ Created PDF: {pdf_path}")
        return pdf_path

    create_pdf()
except ImportError:
    print("! reportlab not installed - skipping PDF")

# Try to create a simple PowerPoint
try:
    from pptx import Presentation

    def create_pptx():
        pptx_path = Path(__file__).parent / "test_presentation.pptx"
        prs = Presentation()

        # Slide 1: Title
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Test LinkedIn Presentation"
        subtitle.text = "Created by chuk-mcp-linkedin"

        # Slide 2: Content
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = "Key Points"
        tf = content.text_frame
        tf.text = "Point 1: Test document upload"
        p = tf.add_paragraph()
        p.text = "Point 2: Verify LinkedIn displays it"
        p.level = 0

        # Slide 3: More content
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = "Conclusion"
        tf = content.text_frame
        tf.text = "LinkedIn document API test successful!"

        prs.save(str(pptx_path))
        print(f"✓ Created PPTX: {pptx_path}")
        return pptx_path

    create_pptx()
except ImportError:
    print("! python-pptx not installed - skipping PPTX")

# Try to create a simple Word document
try:
    from docx import Document

    def create_docx():
        docx_path = Path(__file__).parent / "test_document.docx"
        doc = Document()

        # Add title
        doc.add_heading("Test LinkedIn Document", 0)

        # Add paragraphs
        doc.add_paragraph("This is a test Word document for LinkedIn upload.")
        doc.add_paragraph("Created by chuk-mcp-linkedin")

        doc.add_heading("Section 1", level=1)
        doc.add_paragraph(
            "LinkedIn supports uploading documents in various formats including "
            "PDF, PowerPoint, and Word."
        )

        doc.add_heading("Section 2", level=1)
        doc.add_paragraph("This document tests that functionality.")

        # Add a list
        doc.add_paragraph("Key features:", style="List Bullet")
        doc.add_paragraph("Document upload via API", style="List Bullet")
        doc.add_paragraph("Preview in LinkedIn feed", style="List Bullet")
        doc.add_paragraph("Carousel navigation", style="List Bullet")

        doc.save(str(docx_path))
        print(f"✓ Created DOCX: {docx_path}")
        return docx_path

    create_docx()
except ImportError:
    print("! python-docx not installed - skipping DOCX")

print("\n" + "=" * 60)
print("SUMMARY")
print("=" * 60)
print("Run: pip install reportlab python-pptx python-docx")
print("to generate all test files.")
